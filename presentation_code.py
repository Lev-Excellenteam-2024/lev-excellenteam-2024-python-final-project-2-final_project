from pptx import Presentation
from ai_summary import slide_summary_by_AI
import json
import asyncio

# Asynchronous function to read and process a PowerPoint presentation
async def read_presentation(presentation_path: str) -> list:
    # Load the presentation from the given path
    presentation = Presentation(presentation_path)
    tasks = []

    # Iterate through each slide in the presentation
    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has text and clean it
            if hasattr(shape, "text"):
                cleaned_text = shape.text.replace('\t', ' ').strip()
                if cleaned_text:
                    slide_text.append(cleaned_text)

        # If there is text on the slide, create a task to process it asynchronously
        if slide_text:
            task = asyncio.create_task(process_slide(i, ' '.join(slide_text)))
            tasks.append(task)

    # Wait for all tasks to complete and gather results
    return await asyncio.gather(*tasks)


# Asynchronous function to process each slide's text
async def process_slide(slide_number: int, slide_text: str) -> str:
    try:
        # Call the function to summarize slide text using AI
        summary = await slide_summary_by_AI(slide_text)
    except Exception as e:
        # Handle any errors that occur during processing
        summary = f"Error processing slide {slide_number}: {str(e)}"
    return f"Slide {slide_number}:\n{summary}"


# Function to save the processed data to a JSON file
def save_to_json(filename: str, data: list):
    # Save the processed data to a JSON file
    with open(filename, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)
