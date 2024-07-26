import asyncio
import json
import os
from pptx import Presentation
from openai import AsyncOpenAI
import aiofiles

# Ensure you have set the OpenAI API key in your environment variables
# = os.environ.get("OPENAI_API_KEY")

#if not api_key:
#    raise ValueError("OPENAI_API_KEY environment variable not set")

#print(api_key)
#openai = AsyncOpenAI(api_key)
client = AsyncOpenAI(
    # This is the default and can be omitted
    api_key=os.environ.get("OPENAI_API_KEY"),
)
async def send_slide_to_openai(slide_text, slide_number,my_prompt):
    try:
        response = await client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": my_prompt+slide_text,
            }
        ],
        model="gpt-3.5-turbo",
    )
        return {"slide": slide_number, "response": response.choices[0].message.content}
    except Exception as e:
        return {"slide": slide_number, "error": str(e)}

async def process_presentation(presentation_path):
    presentation = Presentation(presentation_path)
    tasks = []
    my_prompt = ""
    for i, slide in enumerate(presentation.slides):
        slide_text = "I am a student, I need your help to understand this slide from a presentation, in your answer, simlply explain it to me, don't mention its a slide, and act as if I haven't read the presentation yet (ie in your answer if the presentation mentions a polygraph for example, in your answer don't say 'this tool' without giving its name): "
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        tasks.append(send_slide_to_openai(slide_text.strip(), i + 1, my_prompt))

    responses = await asyncio.gather(*tasks)
    return responses

async def main(presentation_path, output_path):
    responses = await process_presentation(presentation_path)
    async with aiofiles.open(output_path, "w") as outfile:
        await outfile.write(json.dumps(responses, indent=4))

if __name__ == "__main__":
    presentation_path = input("Enter path of a .pptx file:\n")
    base_name = os.path.basename(presentation_path)
    # Split the base name into name and extension
    name, _ = os.path.splitext(base_name)
    output_path = name + ".json"  # Replace with your desired output file path

    asyncio.run(main(presentation_path, output_path))
